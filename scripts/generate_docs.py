# -*- coding: utf-8 -*-
"""
生成参赛材料：Word 项目介绍 + PPT 答辩演示
==========================================
基于项目代码库与已有报告，生成符合赛事提交标准的文档。

运行：D:\python.exe scripts/generate_docs.py
输出：output/项目介绍.docx  +  output/答辩演示.pptx
"""

import os
import sys
from datetime import datetime

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

OUTPUT_DIR = os.path.join(ROOT, "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# =============================================================================
# 颜色 / 样式常量
# =============================================================================
# 暖色调体系
TERRACOTTA = (0xD4, 0x78, 0x5C)
OLIVE = (0x7D, 0x8C, 0x6E)
NAVY = (0x2D, 0x33, 0x40)
WARM_WHITE = (0xFD, 0xFA, 0xF3)
SAND = (0xEF, 0xE7, 0xD6)
DARK_TEXT = (0x2D, 0x33, 0x40)

# Chinese quotation marks
LQ = '\u201c'  # left double quote
RQ = '\u201d'  # right double quote

# =============================================================================
# 内容数据
# =============================================================================

PROJECT_NAME = '全龄段AI语言教练'
PROJECT_SUBTITLE = '拼音 \u00b7 英语口语 \u00b7 多语种自由切换'
PROJECT_TAGLINE = '面向所有人的AI语言教练'
COMPETITION = '火山杯 Agent 创新大赛'
COMPETITION_TRACK = 'Trae 赛道'
DATE = '2026年7月'

# --- Word 文档内容 ---
WORD_SECTIONS = [
    {
        'title': '一、项目概述',
        'content': [
            ('本智能体是一个面向所有群体的AI语言教练。无论是刚学拼音的孩童、想矫正口音的成年人、准备出国旅行的退休老人、备考的初高中生、还是零基础的完全初学者，都可以使用。', None),
            ('核心定位是将三大语言学习板块整合于单一智能体中：', None),
            ('板块一：中文拼音学习', '面向零基础拼音学习者、方言口音矫正需求者。从声母韵母认读到拼读应用，覆盖前后鼻音、平翘舌等方言痛点。'),
            ('板块二：英语口语学习', '重点偏向口语训练。从自然拼读入门到场景化对话练习，解决' + LQ + '哑巴英语' + RQ + '问题，支持发音纠错与口音矫正。'),
            ('板块三：多语种自由切换', '支持日语、韩语、法语、西班牙语等，语种间可互换互转，各语种学习进度独立保存。'),
            ('', None),
            ('智能体核心能力：实时识别用户水平与年龄特征、自动匹配对应语种与学习流程、自主规划个性化学习路径、动态学情记录与复习、语音与文字双模态交互。', None),
        ]
    },
    {
        'title': '二、智能体核心架构',
        'content': [
            ('技术架构由工具集成层与三层系统提示词两大模块构成。', None),
            ('工具集成层：', None),
            ('  \u2022 语音合成：调用 Edge-TTS 接口，实现多语种文本转语音，支持调节语速适配老人与儿童', None),
            ('  \u2022 实时语音API：接入多语种 TTS，支持英语、日语、韩语、法语等多语言实时语音输出', None),
            ('  \u2022 语音识别：支持多语种语音输入识别，包括发音评测、口音检测与方言矫正', None),
            ('  \u2022 图文工具：生成口语练习素材、声浪对比图、舌位示意图等可视化教学辅助内容', None),
            ('', None),
            ('三层系统提示词架构：', None),
            ('第一层：基础人设层', '定义智能体为' + LQ + '全龄段AI语言教练' + RQ + '，实时识别用户年龄特征、语言基础和口音问题，自动匹配对应语种与学习流程。对于儿童，系统采用更慢的语速和更简单的引导语言；对于老人，系统增大字体提示、重复关键内容；对于有口音问题的用户，系统启动专项矫正流程。'),
            ('第二层：三大教学分支', '采用强制分支判断机制，自动识别语言需求后锁定对应分支，执行独立教学流程。拼音流程：定级测试\u2192个性化学习计划\u2192声母韵母认读\u2192拼读训练\u2192日常应用。英语口语流程：自然拼读定级\u2192场景化对话\u2192发音矫正\u2192口音改善\u2192自动复习。多语种流程：自由选择目标语种\u2192基础入门\u2192日常会话\u2192语种切换。'),
            ('第三层：长期记忆存储规则', '设计独立的记忆分区（拼音记忆、英语口语记忆、各语种独立记忆），每次学习后自动提取薄弱知识点存入对应分区。每次开启新学习时，系统优先读取记忆库进行针对性复习。这是区别于通用大模型的关键差异点。'),
        ]
    },
    {
        'title': '三、六层生产级工作流',
        'content': [
            ('本智能体在代码层面实现了六层端到端工作流，确保每次调用都稳定、可控、可追溯：', None),
            ('', None),
            ('第1层：意图识别与分支路由', '基于关键词匹配 + LLM 语义理解，将用户输入路由到拼音/英语/多语种/兜底四个分支，计算置信度分数。低置信度时主动追问用户。'),
            ('第2层：记忆检索', '从用户记忆文件中读取学情上下文，包括历史薄弱点、学习进度、开场白使用记录，注入到系统提示词中。'),
            ('第3层：Agentic RAG 检索', '对知识库执行多角度检索：原问题 + 子问题拆解 + 关键词提取。召回相关教学语料，按置信度分高/中/低三级，低置信度时降低模型温度、增加约束。'),
            ('第4层：多候选内容生成', '生成 3 个候选回复（温度 0.8, Top-P 0.9），每个候选注入不同的 Verbalized Sampling 策略（角色扮演/类比解释/分步引导）。层2 SRT 语义排斥指令注入，避免重复表达。'),
            ('第5层：质量校验', '四层护栏系统：\u2460 输入验证（注入检测/范围外话题/不当内容）\u2192 \u2461 输出过滤（LLM-as-a-Judge 二次校验）\u2192 \u2462 行为策略（儿童/老人人群策略）\u2192 \u2463 运行时可观测。SVM 分类器对最终输出评分。未通过自动重生成（最多2次）。'),
            ('第6层：记忆写入', '自动提取回复中的薄弱点标记，写入用户记忆文件。记录本次评估数据（耗时、置信度、质量分数），更新 SRT 排斥窗口和指纹去重窗口。'),
            ('', None),
            ('所有模块均具备优雅降级机制：Ollama 不可用时用规则引擎兜底，sentence-transformers 不可用时降级为哈希指纹，RAG 不可用时跳过检索层。', None),
        ]
    },
    {
        'title': '四、全龄段自适应',
        'content': [
            ('本智能体的核心差异化之一是面向所有人。通过实时识别用户年龄特征，自动适配不同群体的教学策略：', None),
            ('', None),
            ('学龄前及小学儿童（3-12岁）', '趣味化引导，语速 0.5-1.2 倍速。拼音以声母韵母认读为主，英语以自然拼读入门。苏格拉底式提问引导孩子自己发现错误。'),
            ('初高中生（13-18岁）', '对接校内知识体系与考试需求。英语口语融入中考/高考口语考试场景。方言区学生前后鼻音/平翘舌专项矫正。'),
            ('大学生及职场人士（19-40岁）', '英语口语聚焦职场沟通、面试模拟、学术讨论。多语种板块满足留学、出差、跨文化社交需求。高密度内容输出。'),
            ('中老年用户（40岁以上）', '适老化配置：大字体提示、慢速发音示范、关键内容重复三遍。实用旅游情景（问路、购物、点餐）。操作流程极简化。'),
            ('口音问题用户（所有年龄段）', '语音识别检测口音类型，自动生成专项矫正计划。声浪对比图可视化发音差异。反复对比练习逐步改善。'),
        ]
    },
    {
        'title': '五、六层防模板化系统',
        'content': [
            ('为让AI回复像人一样自然多变，实现六层防模板化技术：', None),
            ('', None),
            ('层1：动态提示词工程', '从 12 种开场白、8 种语气、6 种教学策略中轮换，避免连续使用相同元素。'),
            ('层2：语义排斥技术（SRT）', '维护滑动窗口（最近20条输出），抽取高频语块生成排斥列表，注入到生成提示词中，杜绝重复句式。'),
            ('层3：多候选生成 + 评选', '生成 3 个候选，按质量（长度/结构/教学元素）\u00d7 新颖度（1-与历史最大相似度）评分，选最优。'),
            ('层4：Verbalized Sampling', '为每个候选注入不同的表达策略指令（角色扮演/类比解释/分步引导），激活模型预训练多样性。'),
            ('层5：内容指纹去重', 'sentence-transformers 语义嵌入或哈希指纹，余弦相似度超过 0.85 即判定重复，自动过滤。'),
            ('层6：RAG 多样性检索', '查询重写 + 多路召回，从知识库中获取多样化的教学素材，避免反复引用同一段内容。'),
        ]
    },
    {
        'title': '六、市场需求分析',
        'content': [
            ('语言学习赛道正在经历AI驱动的结构性变革：', None),
            ('', None),
            ('整体市场规模', '2026年中国在线语言学习软件市场规模预计达 153.4 亿元，同比增长 39.7%。全球语言学习市场前五大参与者合计市场份额仅 14%，市场高度分散。'),
            ('英语口语学习需求', '国内英语口语应用用户规模突破 1.7 亿。72% 的英语学习者最大难题是' + LQ + '缺乏真实对话环境' + RQ + '。2025年语音矫正培训市场突破 120 亿元。'),
            ('拼音学习需求', '全国约 42% 的小学一年级学生声母发音辨识错误，南方方言区错误率高达 58%。智能拼音教学设备渗透率达 68%。'),
            ('AI 冲击传统赛道', '2026年Q1，多邻国日活首次下降 7%，市值蒸发超 12 亿美元。通用大模型正在蚕食传统语言学习App市场份额，但自身缺乏结构化学习体系。'),
            ('', None),
            ('市场正处于独特窗口期：传统App因AI冲击增长放缓，通用大模型因缺乏教学结构无法真正替代系统化学习。谁能把AI的对话能力与结构化的教学体系结合，并覆盖所有人群，谁就能占据下一个增长高地。', None),
        ]
    },
    {
        'title': '七、竞品对比',
        'content': [
            ('对比维度分析：', None),
            ('  \u2022 多语种支持：传统App 单一语种独立 \u2192 通用大模型 无结构化切换 \u2192 本智能体 三板块+多语种自由切换', None),
            ('  \u2022 英语口语：传统App 以词汇/语法为主 \u2192 通用大模型 有陪练无教学闭环 \u2192 本智能体 口语为核心+场景化+矫正', None),
            ('  \u2022 学习路径：传统App 预设固定路径 \u2192 通用大模型 无预设路径 \u2192 本智能体 定级测试\u2192自动生成\u2192动态调整', None),
            ('  \u2022 长期记忆：传统App 有记录无智能复习 \u2192 通用大模型 无跨会话记忆 \u2192 本智能体 独立记忆分区+自动复习', None),
            ('  \u2022 人群覆盖：传统App 定位单一人群 \u2192 通用大模型 同一风格 \u2192 本智能体 全龄段自适应', None),
            ('  \u2022 口音矫正：传统App 表层打分 \u2192 通用大模型 无专项流程 \u2192 本智能体 方言检测\u2192专项矫正\u2192声浪对比', None),
            ('  \u2022 开发成本：传统App 高 \u2192 通用大模型 低但无教学体系 \u2192 本智能体 5-7天快速验证', None),
            ('', None),
            ('核心竞争策略：在' + LQ + '结构化教学体系' + RQ + '与' + LQ + 'AI对话智能' + RQ + '的交叉地带建立差异化优势，通过全龄段适配和多语种自由切换实现传统产品无法覆盖的人群广度。', None),
        ]
    },
    {
        'title': '八、核心优势与差异化壁垒',
        'content': [
            ('优势一：全龄段自适应，一个智能体覆盖所有人', '为3岁儿童降低语速、增加趣味引导；为70岁老人启用大字体、慢速示范、重复教学；为口音用户提供专项矫正。从孩童到老人，一个智能体陪伴终身学习。'),
            ('优势二：三板块+多语种自由切换，市场无直接竞品', '拼音学习App、英语口语App、其他语种App各自独立。本智能体将三大板块整合，第三板块支持多种语言自由切换，各语种进度独立保存。'),
            ('优势三：英语口语为核心，直击' + LQ + '哑巴英语' + RQ + '痛点', '72%的学习者最大难题是' + LQ + '缺乏真实对话环境' + RQ + '。通过场景化对话（餐厅点餐、机场问路、职场沟通），让用户在低压力环境中高频练习口语。'),
            ('优势四：独立记忆分区，实现' + LQ + '越用越懂你' + RQ, '按语种划分独立记忆分区，每次学习后自动提取薄弱知识点。跨天重启后，智能体主动调用之前的薄弱项进行复习。'),
            ('优势五：三层提示词架构，赋予自主规划能力', '具备感知、记忆、规划、行动、反思五大 Agent 核心能力，能够主动规划学习路径、发现用户卡壳、自动调整教学策略。'),
            ('优势六：开发成本低，快速验证可行', '基于成熟平台与工具链，5-7天即可完成开发与验证。相比传统App动辄数月的开发周期，具有极高的成本效率。'),
        ]
    },
    {
        'title': '九、开发实施路径',
        'content': [
            ('第1天：需求梳理与功能定版', '确定三大板块教学范围、全龄段适配逻辑、定级测试、记忆规则。'),
            ('第2-3天：整理知识文档并导入平台', '拼音发音规则与方言矫正数据、英语口语场景库、多语种基础教材。'),
            ('第4-5天：编写系统提示词与调试工具链', '三层提示词（含全龄段适配+多语种切换）、Edge-TTS/语音API/语音识别接入。'),
            ('第6天：录制与调试演示视频', '拼音完整流程、英语口语场景对话、多语种切换、记忆功能跨天测试、老人/儿童场景演示。'),
            ('第7天：全流程测试、修复Bug、制作答辩PPT', '架构图、三大知识库、全龄段适配对比、与普通App及豆包的对比展示。'),
            ('', None),
            ('技术栈：Coze 平台（Bot 人设/知识库/工作流/数据库）+ 豆包大模型 + Edge-TTS + 本地 Python 生产级参考实现（Ollama + Chroma + SVM）。', None),
        ]
    },
    {
        'title': '十、技术亮点总结',
        'content': [
            ('1. 六层端到端工作流：意图识别 \u2192 记忆检索 \u2192 RAG \u2192 内容生成 \u2192 质量校验 \u2192 记忆写入，每一步都有输入约束、输出校验和失败兜底。', None),
            ('2. 四层护栏系统：输入验证（注入检测/范围外话题/不当内容）\u2192 输出过滤（LLM二次校验）\u2192 行为策略（人群策略）\u2192 运行时可观测。', None),
            ('3. 六层防模板化：动态提示词 \u2192 SRT语义排斥 \u2192 多候选评选 \u2192 Verbalized Sampling \u2192 指纹去重 \u2192 RAG多样性检索。', None),
            ('4. 全模块优雅降级：Ollama 不可用\u2192规则引擎；sentence-transformers 不可用\u2192哈希指纹；RAG 不可用\u2192跳过检索层。', None),
            ('5. 阈值全配置驱动：六环节阈值均有数据依据和注释，生产环境可一键调参。', None),
            ('6. 前后端分离架构：Python 引擎纯逻辑 \u2192 Flask API \u2192 Web 交互界面（暖色调反模板化设计）。', None),
        ]
    },
]

# --- PPT 内容 ---
PPT_SLIDES = [
    {
        'type': 'title',
        'title': PROJECT_NAME,
        'subtitle': PROJECT_SUBTITLE + '\n' + COMPETITION + ' \u00b7 ' + COMPETITION_TRACK + '\n' + DATE,
    },
    {
        'type': 'content',
        'title': '项目概述',
        'bullets': [
            '面向所有年龄段、所有基础水平的AI语言教练',
            '三大教学板块：拼音学习 / 英语口语 / 多语种自由切换',
            '覆盖人群：3-12岁儿童 / 13-18岁青少年 / 19-40岁成人 / 40岁以上老人',
            '核心能力：实时识别年龄与水平 \u2192 自动匹配语种与流程 \u2192 自主规划学习路径 \u2192 动态学情记录与复习',
            '解决传统App' + LQ + '有体系无智能' + RQ + ' + 通用大模型' + LQ + '有智能无体系' + RQ + '的双重困境',
        ]
    },
    {
        'type': 'content',
        'title': '核心痛点',
        'bullets': [
            '传统语言学习App痛点：内容碎片化缺乏语境 | 无长期记忆 | 单一语种切换成本高 | 人群定位狭窄 | 纠音精度不足',
            '通用大模型结构性盲区：无结构化学习路径 | 无长期记忆 | 无全龄段适配 | 无法提供系统化教学',
            '市场数据：72%学习者最大难题是' + LQ + '缺乏真实对话环境' + RQ + ' | 多邻国2026Q1日活下降7% | 南方方言区拼音错误率58%',
            '市场窗口：传统App增长放缓 + 通用大模型无法替代系统化学习 = 巨大机会空间',
        ]
    },
    {
        'type': 'content',
        'title': '核心架构：三层提示词',
        'bullets': [
            '第一层：基础人设层 \u2014 定义Agent身份，实时识别年龄/语言/口音，自动适配',
            '第二层：三大教学分支 \u2014 强制分支判断，锁定拼音/英语/多语种独立教学流程',
            '第三层：长期记忆存储 \u2014 独立记忆分区，薄弱点自动提取，跨天主动复习',
            '',
            '工具集成：Edge-TTS语音合成 | 多语种语音识别 | 声浪对比图 | 舌位示意图',
        ]
    },
    {
        'type': 'content',
        'title': '六层生产级工作流',
        'bullets': [
            'L1 意图识别：关键词匹配 + LLM语义理解 \u2192 分支路由 + 置信度计算',
            'L2 记忆检索：读取用户学情上下文（薄弱点/进度/历史）',
            'L3 Agentic RAG：多角度检索（原问题 + 子问题拆解 + 关键词提取）',
            'L4 多候选生成：3候选 + Verbalized Sampling + SRT语义排斥',
            'L5 质量校验：四层护栏（输入\u2192输出\u2192策略\u2192观测）+ SVM评分 + 失败重生成',
            'L6 记忆写入：薄弱点提取 + 评估记录 + 指纹窗口更新',
            '全模块优雅降级：Ollama不可用\u2192规则引擎 | ST不可用\u2192哈希指纹 | RAG不可用\u2192跳过',
        ]
    },
    {
        'type': 'content',
        'title': '全龄段自适应',
        'bullets': [
            '儿童（3-12岁）：趣味化引导 + 0.5-1.2倍速 + 苏格拉底式提问',
            '青少年（13-18岁）：对接考试需求 + 方言区专项矫正 + 游戏化进度',
            '成人（19-40岁）：职场沟通 + 面试模拟 + 学术讨论 + 高密度输出',
            '老人（40岁以上）：大字体 + 慢速示范 + 关键内容重复三遍 + 极简操作',
            '口音用户：方言检测 \u2192 专项矫正 \u2192 声浪对比图 \u2192 反复练习',
            '一个智能体，从孩童到老人，陪伴终身学习',
        ]
    },
    {
        'type': 'content',
        'title': '六层防模板化系统',
        'bullets': [
            'L1 动态提示词：12种开场白 \u00d7 8种语气 \u00d7 6种策略轮换',
            'L2 SRT语义排斥：滑动窗口抽取高频语块 \u2192 注入排斥指令',
            'L3 多候选评选：质量 \u00d7 新颖度 评分，选最优候选',
            'L4 Verbalized Sampling：角色扮演/类比解释/分步引导策略注入',
            'L5 指纹去重：语义嵌入 + 余弦相似度 > 0.85 自动过滤',
            'L6 RAG多样性：查询重写 + 多路召回，避免重复引用',
        ]
    },
    {
        'type': 'content',
        'title': '竞品对比',
        'bullets': [
            '多语种支持：传统App单一语种 \u2192 豆包支持但无结构化 \u2192 本智能体三板块+自由切换',
            '学习路径：多邻国固定路径 \u2192 豆包无预设 \u2192 本智能体定级+自动生成+动态调整',
            '长期记忆：百词斩有记录无智能复习 \u2192 豆包无跨会话记忆 \u2192 本智能体独立记忆+自动复习',
            '人群覆盖：传统App定位单一人群 \u2192 豆包同一风格 \u2192 本智能体全龄段自适应',
            '口音矫正：传统App表层打分 \u2192 豆包无专项流程 \u2192 本智能体检测\u2192矫正\u2192对比',
            '核心策略：在' + LQ + '结构化教学体系' + RQ + '与' + LQ + 'AI对话智能' + RQ + '的交叉地带建立差异化',
        ]
    },
    {
        'type': 'content',
        'title': '市场需求与机会',
        'bullets': [
            '2026年中国在线语言学习软件市场规模 153.4亿元，同比增长 39.7%',
            '英语口语用户规模突破 1.7亿，语音矫正培训市场 120亿元',
            '全国约42%小学生声母发音辨识错误，南方方言区高达58%',
            '多邻国2026Q1日活下降7%，市值蒸发超12亿美元 \u2014 市场洗牌信号',
            '全球语言学习市场前五大参与者合计份额仅14%，高度分散',
            '核心机会：AI对话能力 \u00d7 结构化教学体系 \u00d7 全人群覆盖',
        ]
    },
    {
        'type': 'content',
        'title': '核心优势',
        'bullets': [
            '全龄段自适应：一个智能体从3岁到70岁，陪伴终身学习',
            '三板块+多语种切换：市场无直接竞品，各语种独立记忆',
            '英语口语为核心：场景化对话 + 发音矫正 + 口音改善',
            '独立记忆分区：越用越懂你，跨天主动复习薄弱点',
            '三层提示词架构：赋予Agent感知/记忆/规划/行动/反思能力',
            '低开发成本：5-7天快速验证，基于成熟平台与工具链',
        ]
    },
    {
        'type': 'content',
        'title': '开发实施路径',
        'bullets': [
            'Day 1：需求梳理与功能定版',
            'Day 2-3：整理知识文档并导入平台',
            'Day 4-5：编写系统提示词与调试工具链',
            'Day 6：录制与调试演示视频',
            'Day 7：全流程测试、修复Bug、制作答辩材料',
            '',
            '技术栈：Coze平台 + 豆包大模型 + Edge-TTS + Python(Ollama+Chroma+SVM)',
        ]
    },
    {
        'type': 'content',
        'title': '总结',
        'bullets': [
            '解决传统App' + LQ + '有体系无智能' + RQ + ' + 通用大模型' + LQ + '有智能无体系' + RQ + '的双重困境',
            '三大板块 + 多语种自由切换，市场无直接竞品',
            '全龄段自适应，一个智能体覆盖3-70岁所有人群',
            '六层工作流 + 四层护栏 + 六层防模板化，生产级品质',
            '长期记忆 + 薄弱点复习，' + LQ + '越用越懂你' + RQ + '的个性化体验',
            '5-7天快速验证，高成本效率，市场窗口期最佳切入时机',
        ]
    },
]


# =============================================================================
# Word 文档生成
# =============================================================================
def generate_word():
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # 页面设置
    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    # 样式
    style = doc.styles['Normal']
    style.font.name = 'Microsoft YaHei'
    style.font.size = Pt(11)
    style.font.color.rgb = RGBColor(*DARK_TEXT)
    style.paragraph_format.line_spacing = 1.5
    style.paragraph_format.space_after = Pt(6)
    style.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    # ===== 封面 =====
    for _ in range(6):
        doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(PROJECT_NAME)
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*TERRACOTTA)
    run.font.name = 'Microsoft YaHei'
    run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(PROJECT_SUBTITLE)
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(*OLIVE)
    run.font.name = 'Microsoft YaHei'
    run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(COMPETITION + ' \u00b7 ' + COMPETITION_TRACK)
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(*NAVY)
    run.font.name = 'Microsoft YaHei'
    run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(DATE)
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.name = 'Microsoft YaHei'
    run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    doc.add_page_break()

    # ===== 目录 =====
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run('\u76ee  \u5f55')
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*NAVY)
    run.font.name = 'Microsoft YaHei'
    run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    doc.add_paragraph()
    for sec in WORD_SECTIONS:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(8)
        run = p.add_run(sec['title'])
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(*OLIVE)
        run.font.name = 'Microsoft YaHei'
        run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    doc.add_page_break()

    # ===== 正文 =====
    for sec in WORD_SECTIONS:
        # 标题
        p = doc.add_paragraph()
        run = p.add_run(sec['title'])
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*TERRACOTTA)
        run.font.name = 'Microsoft YaHei'
        run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after = Pt(8)

        # 标题下划线
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement('w:pBdr')
        bottom = OxmlElement('w:bottom')
        bottom.set(qn('w:val'), 'single')
        bottom.set(qn('w:sz'), '6')
        bottom.set(qn('w:space'), '4')
        bottom.set(qn('w:color'), 'D4785C')
        pBdr.append(bottom)
        pPr.append(pBdr)

        # 内容
        for item in sec['content']:
            text, sub = item if isinstance(item, tuple) else (item, None)

            if not text and not sub:
                doc.add_paragraph()
                continue

            if sub:
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                run = p.add_run(text + '\uff1a')
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*NAVY)
                run.font.name = 'Microsoft YaHei'
                run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

                run = p.add_run(sub)
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*DARK_TEXT)
                run.font.name = 'Microsoft YaHei'
                run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
            else:
                p = doc.add_paragraph()
                p.paragraph_format.first_line_indent = Cm(0.74)
                run = p.add_run(text)
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*DARK_TEXT)
                run.font.name = 'Microsoft YaHei'
                run.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    path = os.path.join(OUTPUT_DIR, '\u9879\u76ee\u4ecb\u7ecd.docx')
    doc.save(path)
    print(f'  \u2705 Word\u6587\u6863\u5df2\u751f\u6210: {path}')
    return path


# =============================================================================
# PPT 生成
# =============================================================================
def generate_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for i, slide_data in enumerate(PPT_SLIDES):
        slide = prs.slides.add_slide(blank_layout)

        # 背景
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*WARM_WHITE)

        # 左侧装饰色条
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = RGBColor(*TERRACOTTA)
        left_bar.line.fill.background()

        # 底部装饰条
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(7.2), Inches(13.18), Inches(0.3)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = RGBColor(*SAND)
        bottom_bar.line.fill.background()

        if slide_data['type'] == 'title':
            # ===== 封面 =====
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0), Inches(13.18), Inches(7.2)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(*WARM_WHITE)
            bg_shape.line.fill.background()

            # 装饰圆形
            circle1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(10.5), Inches(-1), Inches(4), Inches(4)
            )
            circle1.fill.solid()
            circle1.fill.fore_color.rgb = RGBColor(*SAND)
            circle1.line.fill.background()

            circle2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3)
            )
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = RGBColor(0xE0, 0xD5, 0xC0)
            circle2.line.fill.background()

            # 主标题
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*TERRACOTTA)
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.LEFT

            # 副标题
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(2.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = slide_data['subtitle']
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(*NAVY)
            p2.font.name = 'Microsoft YaHei'
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = Pt(36)

        else:
            # ===== 内容页 =====
            # 标题
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*TERRACOTTA)
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.LEFT

            # 标题下划线
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*OLIVE)
            line.line.fill.background()

            # 内容
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.2))
            tf = txBox.text_frame
            tf.word_wrap = True

            for j, bullet in enumerate(slide_data['bullets']):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                if not bullet.strip():
                    p.text = ''
                    p.font.size = Pt(10)
                    continue

                # 检查是否是子标题（以\u2014结尾且不含|）
                is_header = bullet.endswith('\u2014') and '|' not in bullet

                if is_header:
                    p.text = '\u25b8 ' + bullet
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(*NAVY)
                else:
                    p.text = '\u2022 ' + bullet
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*DARK_TEXT)

                p.font.name = 'Microsoft YaHei'
                p.space_after = Pt(10)
                p.line_spacing = Pt(28)

        # 页码
        if slide_data['type'] != 'title':
            txBox = slide.shapes.add_textbox(Inches(12), Inches(7.22), Inches(1), Inches(0.25))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f'{i + 1} / {len(PPT_SLIDES)}'
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.RIGHT

    path = os.path.join(OUTPUT_DIR, '\u7b54\u8fa9\u6f14\u793a.pptx')
    prs.save(path)
    print(f'  \u2705 PPT\u5df2\u751f\u6210: {path}')
    return path


# =============================================================================
# 主入口
# =============================================================================
if __name__ == '__main__':
    print("""
\u2554\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2557
\u2551     \u53c2\u8d5b\u6750\u6599\u751f\u6210\u5668                                         \u2551
\u2551     \u751f\u6210 Word \u9879\u76ee\u4ecb\u7ecd + PPT \u7b54\u8fa9\u6f14\u793a                         \u2551
\u255a\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u255d
""")
    w = generate_word()
    p = generate_ppt()
    print(f'\n  \U0001f4c1 \u8f93\u51fa\u76ee\u5f55: {OUTPUT_DIR}')
    print(f'  \U0001f4c4 Word: {os.path.basename(w)}')
    print(f'  \U0001f4ca PPT:  {os.path.basename(p)}')
    print('\n  \u751f\u6210\u5b8c\u6210\uff01')